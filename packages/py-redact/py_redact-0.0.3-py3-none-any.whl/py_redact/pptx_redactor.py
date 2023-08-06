import re
import logging

from pptx import Presentation

logger = logging.getLogger("pptx-redact")


class PptxRedactor:
    def __init__(self, ppt_obj_path, regexes, replace_char):
        self.ppt_obj_path = ppt_obj_path
        self.regexes = regexes
        self.replace_char = replace_char

    def __redact_helper__(self, shapes, notes_slides):
        """
         Helper function for the redact function
        :param shapes (list[shape]): list containing all shapes in all slides of the presentation
        :param notes_slides (list[notes_slide]): list containing all notes_slides associated with every slide
        :return:
        """
        for reg in self.regexes:
            regex = re.compile(reg)
            for shape in shapes:
                if shape.has_text_frame:
                    if regex.search(shape.text):
                        text_frame = shape.text_frame
                        for paragraph in text_frame.paragraphs:
                            whole_text = "".join(run.text for run in paragraph.runs)
                            inline = paragraph.runs
                            for i in range(len(inline)):
                                if regex.search(inline[i].text):
                                    text = regex.sub(self.replace_char * len(inline[i].text), inline[i].text)
                                    inline[i].text = text
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if (regex.search(cell.text)):
                                text = regex.sub(self.replace_char * len(cell.text), cell.text)
                                cell.text = text
            for notes_slide in notes_slides:
                text_frame = notes_slide.notes_text_frame
                m = regex.search(text_frame.text)
                if (m):
                    text = regex.sub(self.replace_char * len(text_frame.text), text_frame.text)
                    text_frame.text = text

    def redact(self, output_file_path):
        """
        Redacts the given .pptx file and writes result to output file
        :param output_file_path (string): path of the file to write the result
        :return:
        """
        help(self.__redact_helper__)
        prs = Presentation(self.ppt_obj_path)
        # To get shapes in your slides
        slides = [slide for slide in prs.slides]
        notes_slides = []
        shapes = []
        for slide in slides:
            for shape in slide.shapes:
                shapes.append(shape)
            if slide.has_notes_slide:
                notes_slides.append(slide.notes_slide)
        self.__redact_helper__(shapes, notes_slides)
        if self.ppt_obj_path == output_file_path:
            logger.warning("Input and Output files are same!")
        logger.info(" Updated file saved as: " + output_file_path)
        prs.save(output_file_path)
